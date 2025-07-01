import os
import pikepdf
import tempfile
from flask import Flask, request, send_file, render_template
from io import BytesIO
from PyPDF2 import PdfReader, PdfWriter
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.colors import Color
from pdf2image import convert_from_bytes
from pptx import Presentation
from pptx.util import Inches
from docx import Document
from PIL import Image, ImageDraw, ImageFont, ImageEnhance
import zipfile
import pythoncom


# Import for Windows COM automation for PPT to PDF and Word to PDF
try:
    import comtypes.client
    COMTYPES_AVAILABLE = True
except ImportError:
    COMTYPES_AVAILABLE = False

app = Flask(__name__)

def pdf_contains_images(pdf):
    # Check if any page contains images (XObject of subtype /Image)
    for page in pdf.pages:
        resources = page.get('/Resources', {})
        xobjects = resources.get('/XObject', {})
        for obj in xobjects.values():
            try:
                xobj = obj.get_object()
                if xobj.get('/Subtype') == '/Image':
                    return True
            except Exception:
                continue
    return False

def add_watermark(input_pdf_stream, watermark_text):
    packet = BytesIO()
    page_width, page_height = letter
    c = canvas.Canvas(packet, pagesize=(page_width, page_height))
    font_size = 20
    c.setFont("Helvetica", font_size)
    watermark_color = Color(0.5, 0.5, 0.5, alpha=0.4)
    c.setFillColor(watermark_color)
    diagonal_step = 150

    for x in range(-int(page_height), int(page_width), diagonal_step):
        for y in range(-int(page_height), int(page_height), diagonal_step):
            c.saveState()
            c.translate(x, y)
            c.rotate(45)
            c.drawString(0, 0, watermark_text)
            c.restoreState()
    c.save()
    packet.seek(0)

    watermark_pdf = PdfReader(packet)
    watermark_page = watermark_pdf.pages[0]

    input_pdf = PdfReader(input_pdf_stream)
    output_pdf = PdfWriter()

    for page in input_pdf.pages:
        page.merge_page(watermark_page)
        output_pdf.add_page(page)

    output_stream = BytesIO()
    output_pdf.write(output_stream)
    output_stream.seek(0)
    return output_stream


@app.route('/')
def index():
    return render_template('index.html')


@app.route('/merge', methods=['POST'])
def merge():
    files = request.files.getlist('pdfs')
    if not files or len(files) < 2:
        return 'Upload at least two PDF files to merge.', 400
    output_pdf = PdfWriter()
    try:
        for f in files:
            reader = PdfReader(f.stream)
            for page in reader.pages:
                output_pdf.add_page(page)
        output_stream = BytesIO()
        output_pdf.write(output_stream)
        output_stream.seek(0)
        return send_file(
            output_stream,
            as_attachment=True,
            download_name="merged.pdf",
            mimetype='application/pdf'
        )
    except Exception as e:
        return f"Error merging PDFs: {e}", 500


@app.route('/split', methods=['POST'])
def split():
    file = request.files.get('pdf')
    start = request.form.get('start')
    end = request.form.get('end')
    if not file or not start or not end:
        return 'Please provide PDF file and start and end pages.', 400
    try:
        start = int(start)
        end = int(end)
    except ValueError:
        return 'Start and end pages must be integers.', 400
    if start < 1 or end < start:
        return 'Invalid page range.', 400
    try:
        input_pdf = PdfReader(file.stream)
        output_pdf = PdfWriter()
        total_pages = len(input_pdf.pages)
        if end > total_pages:
            return f'End page exceeds total pages ({total_pages}).', 400
        for i in range(start - 1, end):
            output_pdf.add_page(input_pdf.pages[i])
        output_stream = BytesIO()
        output_pdf.write(output_stream)
        output_stream.seek(0)
        return send_file(
            output_stream,
            as_attachment=True,
            download_name="split.pdf",
            mimetype='application/pdf'
        )
    except Exception as e:
        return f'Error splitting PDF: {e}', 500


@app.route('/compress', methods=['POST'])
def compress():
    file = request.files.get('pdf')
    if not file:
        return "No PDF uploaded", 400
    try:
        input_pdf = BytesIO(file.read())
        pdf = pikepdf.open(input_pdf)
        # Check for images
        if pdf_contains_images(pdf):
            return "Compression only supported for text-based PDFs (no images found).", 400
        output_pdf = BytesIO()
        pdf.save(output_pdf, compress_streams=True, recompress_flate=True)
        output_pdf.seek(0)
        return send_file(
            output_pdf,
            as_attachment=True,
            download_name='compressed.pdf',
            mimetype='application/pdf'
        )
    except Exception as e:
        return f"Compression failed: {e}", 500


@app.route('/watermark', methods=['POST'])
def watermark():
    uploaded_file = request.files.get('pdf')
    watermark_text = request.form.get('watermark_text', '').strip()
    if not uploaded_file or uploaded_file.filename == '':
        return "No PDF file uploaded", 400
    if not watermark_text:
        return "No watermark text provided", 400
    try:
        input_pdf_stream = BytesIO(uploaded_file.read())
        output_stream = add_watermark(input_pdf_stream, watermark_text)
        return send_file(
            output_stream,
            as_attachment=True,
            download_name="watermarked.pdf",
            mimetype="application/pdf"
        )
    except Exception as e:
        return f"Error applying watermark: {e}", 500


@app.route('/image2pdf', methods=['POST'])
def image2pdf():
    files = request.files.getlist('images')
    if not files or len(files) == 0:
        return "No images uploaded", 400

    images = []
    for file in files:
        img = Image.open(file.stream)
        # Convert all images to RGB (PDFs do not support alpha channel)
        if img.mode in ('RGBA', 'P'):
            img = img.convert('RGB')
        images.append(img)

    if not images:
        return "No valid images found", 400

    output_stream = BytesIO()
    # Save all images as a single PDF (first image + rest as pages)
    images[0].save(
        output_stream, format='PDF', save_all=True, append_images=images[1:]
    )
    output_stream.seek(0)
    return send_file(
        output_stream,
        as_attachment=True,
        download_name='converted.pdf',
        mimetype='application/pdf'
    )


@app.route('/pdf2jpeg', methods=['POST'])
def pdf2jpeg():
    file = request.files.get('pdf')
    if not file:
        return "No PDF uploaded", 400
    try:
        images = convert_from_bytes(file.read())
        output_zip = BytesIO()
        with zipfile.ZipFile(output_zip, "w") as zipf:
            for i, image in enumerate(images):
                img_byte_arr = BytesIO()
                image.save(img_byte_arr, format='JPEG')
                zipf.writestr(f'page_{i+1}.jpeg', img_byte_arr.getvalue())
        output_zip.seek(0)
        return send_file(
            output_zip,
            as_attachment=True,
            download_name='pdf_images.zip',
            mimetype='application/zip'
        )
    except Exception as e:
        return f'Error converting PDF to JPEG: {e}', 500


@app.route('/pdf2ppt', methods=['POST'])
def pdf2ppt():
    file = request.files.get('pdf')
    if not file:
        return "No PDF uploaded", 400
    try:
        images = convert_from_bytes(file.read())
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        for image in images:
            slide = prs.slides.add_slide(blank_slide_layout)
            img_byte_arr = BytesIO()
            image.save(img_byte_arr, format='PNG')
            img_byte_arr.seek(0)
            pic = slide.shapes.add_picture(img_byte_arr, Inches(0), Inches(0),
                                           width=prs.slide_width, height=prs.slide_height)
        output_stream = BytesIO()
        prs.save(output_stream)
        output_stream.seek(0)
        return send_file(
            output_stream,
            as_attachment=True,
            download_name='converted.pptx',
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
    except Exception as e:
        return f'Error converting PDF to PPT: {e}', 500


@app.route('/ppt2pdf', methods=['POST'])
def ppt2pdf():
    if not COMTYPES_AVAILABLE:
        return "ppt2pdf requires comtypes package and Windows with MS PowerPoint installed.", 500

    file = request.files.get('ppt')
    if not file:
        return "No PPT uploaded", 400

    with tempfile.TemporaryDirectory() as tmpdir:
        ppt_path = os.path.join(tmpdir, 'input.pptx')
        pdf_path = os.path.join(tmpdir, 'output.pdf')
        file.save(ppt_path)
        try:
            pythoncom.CoInitialize()
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            powerpoint.Visible = 1
            presentation = powerpoint.Presentations.Open(ppt_path, WithWindow=False)
            presentation.SaveAs(pdf_path, FileFormat=32)  # 32 = ppt to pdf
            presentation.Close()
            powerpoint.Quit()

            # Read PDF to memory before temp dir is cleaned up
            with open(pdf_path, 'rb') as f:
                pdf_bytes = f.read()
            mem_pdf = BytesIO(pdf_bytes)
            mem_pdf.seek(0)
            return send_file(mem_pdf, as_attachment=True, download_name='converted.pdf', mimetype='application/pdf')
        except Exception as e:
            return f"Error converting PPT to PDF: {e}", 500
        finally:
            pythoncom.CoUninitialize()

   


@app.route('/pdf2word', methods=['POST'])
def pdf2word():
    file = request.files.get('pdf')
    if not file:
        return "No PDF uploaded", 400
    try:
        reader = PdfReader(file.stream)
        doc = Document()
        for page in reader.pages:
            text = page.extract_text()
            if text:
                doc.add_paragraph(text)
        output_stream = BytesIO()
        doc.save(output_stream)
        output_stream.seek(0)
        return send_file(
            output_stream,
            as_attachment=True,
            download_name='converted.docx',
            mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        )
    except Exception as e:
        return f'Error converting PDF to Word: {e}', 500


@app.route('/word2pdf', methods=['POST'])
def word2pdf():
    if not COMTYPES_AVAILABLE:
        return "word2pdf requires comtypes package and Windows with MS Word installed.", 500

    file = request.files.get('word')
    if not file:
        return "No Word file uploaded", 400

    with tempfile.TemporaryDirectory() as tmpdir:
        input_docx_path = os.path.join(tmpdir, 'input.docx')
        output_pdf_path = os.path.join(tmpdir, 'output.pdf')
        file.save(input_docx_path)
        try:
            pythoncom.CoInitialize()  # Initialize COM
            word = comtypes.client.CreateObject('Word.Application')
            word.Visible = False
            doc = word.Documents.Open(input_docx_path)
            doc.SaveAs(output_pdf_path, FileFormat=17)  # 17 = doc to pdf
            doc.Close()
            word.Quit()
            
            # Read the file into memory before closing the temp directory
            with open(output_pdf_path, 'rb') as f:
                pdf_data = f.read()
                
            return send_file(
                BytesIO(pdf_data),
                as_attachment=True, 
                download_name='converted.pdf', 
                mimetype='application/pdf'
            )
        except Exception as e:
            return f"Error converting Word to PDF: {e}", 500
        finally:
            pythoncom.CoUninitialize()  # Uninitialize COM




if __name__ == '__main__':
    app.run(debug=True) 
